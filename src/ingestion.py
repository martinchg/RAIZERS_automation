"""
Text_extraction_light.py : Version légère sans GPU.
- HTML  : BeautifulSoup
- DOCX  : python-docx + fallback XML natif
- EXCEL : Pandas
- PDF   : pymupdf4llm  (PyMuPDF natif → markdown propre, pas de modèle ML)
- AUDIO : Whisper (CPU)
- TXT / MD : lecture directe

Interface identique à Text_extraction.py :
    extract(file_path) -> Tuple[List[Dict], Dict[str, int]]

Usage dans main.py — remplacer simplement l'import :
    from Text_extraction_light import extract
"""

import logging
import warnings
from pathlib import Path
from typing import Any, Dict, List, Tuple
from zipfile import ZipFile
import xml.etree.ElementTree as ET

import fitz  # PyMuPDF (déjà dans requirements)

# --- Imports dynamiques ---
try:
    import pymupdf4llm
except ImportError:
    pymupdf4llm = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import pandas as pd
except ImportError:
    pd = None


# --- Configuration ---
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
warnings.filterwarnings("ignore")

_WORD_XML_NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
_PPT_XML_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


def _append_unique(parts: List[str], text: str) -> None:
    cleaned = " ".join((text or "").split()).strip()
    if cleaned:
        parts.append(cleaned)


def _extract_docx_xml_fallback(file_path: Path) -> List[str]:
    """Fallback stdlib pour DOCX si python-docx échoue sur Streamlit."""
    xml_parts: List[str] = []
    try:
        with ZipFile(file_path) as archive:
            names = sorted(
                name
                for name in archive.namelist()
                if name == "word/document.xml"
                or name.startswith("word/header")
                or name.startswith("word/footer")
            )
            for name in names:
                root = ET.fromstring(archive.read(name))
                for paragraph in root.findall(".//w:p", _WORD_XML_NS):
                    text_nodes = paragraph.findall(".//w:t", _WORD_XML_NS)
                    text = "".join(node.text or "" for node in text_nodes).strip()
                    _append_unique(xml_parts, text)
    except Exception as exc:
        logger.error(f"❌ Fallback XML DOCX impossible ({file_path.name}) : {exc}")
    return xml_parts


def _extract_pptx_xml_fallback(file_path: Path) -> List[str]:
    """Fallback stdlib pour PPTX si python-pptx n'est pas disponible."""
    slides: List[str] = []
    try:
        with ZipFile(file_path) as archive:
            slide_names = sorted(
                name
                for name in archive.namelist()
                if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            )
            for name in slide_names:
                root = ET.fromstring(archive.read(name))
                paragraphs: List[str] = []
                for paragraph in root.findall(".//a:p", _PPT_XML_NS):
                    text = "".join(
                        node.text or ""
                        for node in paragraph.findall(".//a:t", _PPT_XML_NS)
                    ).strip()
                    _append_unique(paragraphs, text)
                slide_text = "\n".join(paragraphs).strip()
                if slide_text:
                    slides.append(slide_text)
    except Exception as exc:
        logger.error(f"❌ Fallback XML PPTX impossible ({file_path.name}) : {exc}")
    return slides



# ---------------------------------------------------------------------------
# 1. HTML (BeautifulSoup) — identique à Text_extraction.py
# ---------------------------------------------------------------------------
def extract_html(file_path: Path) -> List[Dict[str, Any]]:
    if BeautifulSoup is None:
        return []
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")

        # Sauvetage maths LaTeX inline / display
        for script in soup.find_all("script", type="math/tex"):
            tex = script.get_text()
            if "mode=display" in str(script):
                tag = soup.new_tag("p")
                tag.string = f" $$ {tex} $$ "
                script.replace_with(tag)
            else:
                script.replace_with(f" $ {tex} $ ")

        # Sauvetage alt des images
        for img in soup.find_all("img"):
            alt = img.get("alt", "").strip()
            if alt:
                img.replace_with(f" [IMAGE: {alt}] ")
            else:
                img.decompose()

        # Nettoyage balises inutiles
        for tag in soup(["script", "style", "header", "footer", "nav", "meta", "noscript", "iframe"]):
            tag.extract()

        text = soup.get_text(separator="\n").strip()
        clean_text = "\n".join(line.strip() for line in text.splitlines() if line.strip())

        if clean_text:
            return [{"text": clean_text, "category": "HTML", "metadata": {"filename": file_path.name}}]
    except Exception as e:
        logger.error(f"❌ Erreur HTML : {e}")
    return []


# ---------------------------------------------------------------------------
# 2. DOCX (python-docx)
# ---------------------------------------------------------------------------
def extract_docx(file_path: Path) -> List[Dict[str, Any]]:
    parts: List[str] = []
    try:
        from docx import Document
        from docx.document import Document as _Document
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table, _Cell
        from docx.text.paragraph import Paragraph

        def iter_blocks(container):
            if isinstance(container, _Document):
                parent_elm = container.element.body
            elif isinstance(container, _Cell):
                parent_elm = container._tc
            else:
                parent_elm = getattr(container, "_element", None)
            if parent_elm is None:
                return
            for child in parent_elm.iterchildren():
                if isinstance(child, CT_P):
                    yield Paragraph(child, container)
                elif isinstance(child, CT_Tbl):
                    yield Table(child, container)

        def table_to_text(table: Table) -> str:
            rows: List[str] = []
            for row in table.rows:
                cells: List[str] = []
                for cell in row.cells:
                    cell_parts: List[str] = []
                    for block in iter_blocks(cell):
                        if isinstance(block, Paragraph):
                            _append_unique(cell_parts, block.text)
                        elif isinstance(block, Table):
                            nested = table_to_text(block)
                            if nested:
                                cell_parts.append(nested.replace("\n", " ; "))
                    cell_text = " ".join(cell_parts).strip()
                    if cell_text:
                        cells.append(cell_text)
                if cells:
                    rows.append(" | ".join(cells))
            return "\n".join(rows)

        doc = Document(str(file_path))

        for block in iter_blocks(doc):
            if isinstance(block, Paragraph):
                _append_unique(parts, block.text)
            elif isinstance(block, Table):
                table_text = table_to_text(block)
                if table_text:
                    parts.append(table_text)

        seen_headers = set()
        for section in doc.sections:
            for label in ("header", "footer"):
                container = getattr(section, label, None)
                if container is None:
                    continue
                for block in iter_blocks(container):
                    if isinstance(block, Paragraph):
                        text = block.text
                    elif isinstance(block, Table):
                        text = table_to_text(block)
                    else:
                        text = ""
                    cleaned = " ".join(text.split()).strip()
                    if cleaned and cleaned not in seen_headers:
                        seen_headers.add(cleaned)
                        parts.append(f"[{label.upper()}] {cleaned}")
    except Exception as e:
        logger.warning(f"⚠️ DOCX via python-docx en échec ({file_path.name}) : {e}")

    if not parts:
        parts = _extract_docx_xml_fallback(file_path)

    final_text = "\n\n".join(parts)
    if final_text:
        return [{"text": final_text, "category": "DOCX", "metadata": {"filename": file_path.name}}]
    return []

# ---------------------------------------------------------------------------
# 3. PPTX (python-pptx)
# ---------------------------------------------------------------------------

def extract_pptx(file_path: Path) -> List[Dict[str, Any]]:
    parts: List[str] = []
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        for slide in prs.slides:
            slide_parts: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_parts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_parts.append(row_text)
            if slide_parts:
                parts.append("\n".join(slide_parts))
    except Exception as e:
        logger.warning(f"⚠️ PPTX via python-pptx en échec ({file_path.name}) : {e}")

    if not parts:
        parts = _extract_pptx_xml_fallback(file_path)

    final_text = "\n\n".join(parts)
    if final_text:
        return [{"text": final_text, "category": "PPTX", "metadata": {"filename": file_path.name}}]
    return []
# ---------------------------------------------------------------------------
# 4. EXCEL (Pandas) — identique à Text_extraction.py
# ---------------------------------------------------------------------------
def extract_excel(file_path: Path) -> List[Dict[str, Any]]:
    if pd is None:
        return []
    elements: List[Dict[str, Any]] = []
    try:
        excel = pd.ExcelFile(file_path)
        for sheet in excel.sheet_names:
            df = pd.read_excel(excel, sheet_name=sheet)
            df = df.dropna(how="all").dropna(axis=1, how="all")
            if not df.empty:
                text = f"## Feuille: {sheet}\n\n" + df.to_markdown(index=False)
                elements.append({
                    "text": text,
                    "category": "Spreadsheet",
                    "metadata": {"filename": file_path.name, "sheet_name": sheet},
                })
    except Exception as e:
        logger.error(f"❌ Erreur Excel : {e}")
    return elements


# ---------------------------------------------------------------------------
# 5. PDF (pymupdf4llm) — version légère, pas de GPU/modèle ML
# ---------------------------------------------------------------------------
def extract_pdf(file_path: Path) -> List[Dict[str, Any]]:
    """
    Extraction PDF via pymupdf4llm — page par page.
    Chaque page devient un élément séparé avec son numéro de page en metadata.

    ⚠️  Ne convient pas aux PDFs entièrement scannés (images sans couche texte).
        Pour ces cas, utiliser Text_extraction.py (Marker avec OCR).
    """
    if pymupdf4llm is None:
        logger.error("❌ pymupdf4llm n'est pas installé. Lancez : pip install pymupdf4llm")
        return []
    try:
        # Extraction page par page → liste de dicts avec metadata par page
        page_chunks = pymupdf4llm.to_markdown(
            str(file_path),
            show_progress=False,
            page_chunks=True,  # ← retourne une liste de dicts par page
        )

        elements: List[Dict[str, Any]] = []
        for chunk in page_chunks:
            text = chunk.get("text", "").strip()
            if not text:
                continue
            page_num = chunk.get("metadata", {}).get("page_number", None)
            elements.append({
                "text": text,
                "category": "PDF",
                "metadata": {
                    "filename": file_path.name,
                    "page_label": page_num,  # pymupdf4llm numérote déjà à partir de 1
                },
            })

        logger.info(f"  📄 {file_path.name}: {len(elements)} page(s) extraite(s)")
        return elements

    except Exception as e:
        logger.error(f"❌ Erreur PDF ({file_path.name}): {e}")
    return []


# ---------------------------------------------------------------------------
# ROUTER — interface identique à Text_extraction.py
# ---------------------------------------------------------------------------
def extract(file_path: "str | Path", source_path: str = "") -> Tuple[List[Dict], Dict[str, int]]:
    """
    Point d'entrée unique.
    Retourne (elements, stats).

    source_path : chemin relatif dans le projet Dropbox (ex: "2. audit/CG études.pdf")
                  Injecté dans la metadata de chaque élément pour traçabilité.
    """
    path = Path(file_path)
    if not path.exists():
        logger.warning(f"⚠️  Fichier introuvable : {file_path}")
        return [], {}

    ext = path.suffix.lower()

    if ext == ".pdf":
        elements, stats = extract_pdf(path), {"PDF": 1}
    elif ext in {".html", ".htm"}:
        elements, stats = extract_html(path), {"HTML": 1}
    elif ext == ".docx":
        elements, stats = extract_docx(path), {"DOCX": 1}
    elif ext in {".xlsx", ".xls"}:
        elements, stats = extract_excel(path), {"Excel": 1}
    elif ext in {".pptx", ".ppt"}:
        elements, stats = extract_pptx(path), {"PPTX": 1}
    elif ext in {".txt", ".md"}:
        with open(path, "r", encoding="utf-8") as f:
            elements = [{"text": f.read(), "category": "Text", "metadata": {"filename": path.name}}]
        stats = {"Text": 1}
    else:
        logger.warning(f"⚠️  Extension non supportée : {ext}")
        return [], {}

    # Injecter source_path + file_type dans toutes les metadata
    file_type = ext.lstrip(".")  # "pdf", "docx", "xlsx"...
    for el in elements:
        meta = el.setdefault("metadata", {})
        if source_path:
            meta["source_path"] = source_path
        meta["file_type"] = file_type

    return elements, stats
